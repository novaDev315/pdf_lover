#!/usr/bin/env python3
"""Bounded document conversion, OCR, and raster-compression engine."""

from __future__ import annotations

import argparse
import json
import re
import subprocess
from pathlib import Path

MAX_PAGES = 250


def command(args: list[str], *, text: bool = True) -> subprocess.CompletedProcess:
    return subprocess.run(args, check=True, capture_output=True, text=text)


def pdf_page_count(source: Path) -> int:
    result = command(["pdfinfo", str(source)])
    match = re.search(r"^Pages:\s+(\d+)\s*$", result.stdout, re.MULTILINE)
    if not match:
        raise ValueError("Could not determine the PDF page count")
    pages = int(match.group(1))
    if pages < 1 or pages > MAX_PAGES:
        raise ValueError(f"PDF must contain between 1 and {MAX_PAGES} pages")
    return pages


def render_pages(source: Path, workdir: Path, dpi: int, image_format: str) -> list[Path]:
    prefix = workdir / "page"
    command([
        "pdftocairo",
        f"-{image_format}",
        "-r",
        str(dpi),
        str(source),
        str(prefix),
    ])
    extension = "jpg" if image_format == "jpeg" else image_format
    pages = sorted(workdir.glob(f"page-*.{extension}"))
    if not pages:
        raise ValueError("PDF rendering produced no pages")
    return pages


def extract_layout_text(source: Path) -> list[str]:
    result = command(["pdftotext", "-layout", str(source), "-"])
    pages = result.stdout.split("\f")
    if pages and not pages[-1].strip():
        pages.pop()
    return pages


def convert_docx(source: Path, output: Path) -> None:
    from docx import Document

    pages = extract_layout_text(source)
    if not any(page.strip() for page in pages):
        raise ValueError("DOCX conversion found no extractable text; run OCR first")
    document = Document()
    for page_index, page_text in enumerate(pages):
        for paragraph in re.split(r"\n\s*\n", page_text.strip()):
            cleaned = "\n".join(line.rstrip() for line in paragraph.splitlines()).strip()
            if cleaned:
                document.add_paragraph(cleaned)
        if page_index + 1 < len(pages):
            document.add_page_break()
    document.save(output)


def table_rows(page_text: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in page_text.splitlines():
        cells = [cell.strip() for cell in re.split(r"\s{2,}|\t+", line.strip()) if cell.strip()]
        if len(cells) >= 2:
            rows.append(cells)
    return rows


def convert_xlsx(source: Path, output: Path) -> None:
    from openpyxl import Workbook

    workbook = Workbook()
    workbook.remove(workbook.active)
    table_count = 0
    for page_number, page_text in enumerate(extract_layout_text(source), start=1):
        rows = table_rows(page_text)
        if len(rows) < 2:
            continue
        worksheet = workbook.create_sheet(title=f"Page {page_number}")
        for row in rows:
            worksheet.append(row)
        table_count += 1
    if table_count == 0:
        raise ValueError("XLSX conversion found no table-like rows")
    workbook.save(output)


def convert_pptx(source: Path, output: Path, workdir: Path, dpi: int) -> None:
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    images = render_pages(source, workdir, dpi, "png")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    for image_path in images:
        with Image.open(image_path) as image:
            image_ratio = image.width / image.height
        slide_ratio = presentation.slide_width / presentation.slide_height
        if image_ratio > slide_ratio:
            width = presentation.slide_width
            height = int(width / image_ratio)
            left = 0
            top = int((presentation.slide_height - height) / 2)
        else:
            height = presentation.slide_height
            width = int(height * image_ratio)
            left = int((presentation.slide_width - width) / 2)
            top = 0
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    presentation.save(output)


def enhance_image(source: Path, output: Path) -> Path:
    from PIL import Image, ImageEnhance, ImageOps

    with Image.open(source) as image:
        gray = ImageOps.grayscale(image)
        corrected = ImageOps.autocontrast(gray, cutoff=1)
        enhanced = ImageEnhance.Contrast(corrected).enhance(1.2)
        enhanced.save(output, format="PNG", optimize=True)
    return output


def searchable_ocr(
    source: Path,
    output: Path,
    workdir: Path,
    dpi: int,
    language: str,
    enhance: bool,
) -> None:
    images = render_pages(source, workdir, dpi, "png")
    page_pdfs: list[Path] = []
    for index, image in enumerate(images, start=1):
        ocr_image = enhance_image(image, workdir / f"enhanced-{index}.png") if enhance else image
        base = workdir / f"ocr-{index}"
        command(["tesseract", str(ocr_image), str(base), "-l", language, "pdf"])
        page_pdf = base.with_suffix(".pdf")
        if not page_pdf.exists():
            raise ValueError(f"OCR produced no PDF for page {index}")
        page_pdfs.append(page_pdf)
    command(["pdfunite", *[str(page) for page in page_pdfs], str(output)])


def lossy_compress(source: Path, output: Path, workdir: Path, dpi: int, quality: int) -> None:
    from PIL import Image

    images = render_pages(source, workdir, dpi, "jpeg")
    opened = [Image.open(path).convert("RGB") for path in images]
    try:
        opened[0].save(
            output,
            format="PDF",
            save_all=True,
            append_images=opened[1:],
            resolution=dpi,
            quality=quality,
            optimize=True,
        )
    finally:
        for image in opened:
            image.close()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--operation", required=True)
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    parser.add_argument("--options", required=True, type=Path)
    parser.add_argument("--workdir", required=True, type=Path)
    args = parser.parse_args()

    options = json.loads(args.options.read_text(encoding="utf-8"))
    args.workdir.mkdir(mode=0o700, parents=False, exist_ok=False)
    pdf_page_count(args.input)

    if args.operation == "pdf.convert.docx":
        convert_docx(args.input, args.output)
    elif args.operation == "pdf.convert.xlsx":
        convert_xlsx(args.input, args.output)
    elif args.operation == "pdf.convert.pptx":
        convert_pptx(args.input, args.output, args.workdir, options["dpi"])
    elif args.operation == "pdf.ocr":
        searchable_ocr(
            args.input,
            args.output,
            args.workdir,
            options["dpi"],
            options["language"],
            options["enhanceScans"],
        )
    elif args.operation == "pdf.compress.lossy":
        lossy_compress(
            args.input,
            args.output,
            args.workdir,
            options["dpi"],
            options["quality"],
        )
    else:
        raise ValueError(f"Unsupported document operation: {args.operation}")


if __name__ == "__main__":
    try:
        main()
    except (ValueError, json.JSONDecodeError, subprocess.CalledProcessError) as error:
        message = error.stderr.strip() if isinstance(error, subprocess.CalledProcessError) and error.stderr else str(error)
        raise SystemExit(message)
